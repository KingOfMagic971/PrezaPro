# meta_developer: @k1sIotaa
# scope: hijacker
# requires: python-pptx fpdf openai

import os
import json
import asyncio
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from fpdf import FPDF
import openai

# --- НАСТРОЙКИ ---
# Вставь сюда свой ключ или передавай его при инициализации
OPENAI_API_KEY = "sk-..." 

class PresentationGenius:
    """
    Мощный ИИ-модуль для создания презентаций.
    Клон функционала @prezaprobot.
    """
    def __init__(self, api_key=None):
        self.api_key = api_key or OPENAI_API_KEY
        self.client = openai.AsyncOpenAI(api_key=self.api_key)
        
        # Настройки цветов для тем (RGB)
        self.design_palettes = {
            "Ash": {"bg": (240, 240, 240), "text": (50, 50, 50), "accent": (100, 100, 100)},
            "Coal": {"bg": (30, 30, 30), "text": (255, 255, 255), "accent": (255, 215, 0)}, # Черный + Золото
            "Icebreaker": {"bg": (220, 240, 255), "text": (0, 50, 100), "accent": (0, 150, 255)},
            "Electric": {"bg": (10, 10, 40), "text": (200, 255, 255), "accent": (0, 255, 0)},
            "Basic Light": {"bg": (255, 255, 255), "text": (0, 0, 0), "accent": (0, 112, 192)},
        }

    async def generate_structure(self, topic, lang, slides_count, audience, volume, design, img_type):
        """
        Генерирует структуру презентации через GPT-4o.
        Возвращает чистый JSON.
        """
        
        # Сложный системный промпт
        system_prompt = (
            f"Ты — профессиональный создатель презентаций. Твоя задача — создать структуру презентации."
            f"\nЯзык: {lang}"
            f"\nАудитория: {audience}"
            f"\nОбъем текста на слайд: {volume}"
            f"\nКоличество слайдов: {slides_count}"
            f"\nОтвет верни СТРОГО в формате JSON: "
            f'{{"slides": [{{"title": "Заголовок", "content": "Текст", "img_prompt": "Описание картинки"}}]}}'
        )

        user_prompt = f"Тема презентации: {topic}. Стиль дизайна: {design}. Тип картинок: {img_type}."

        try:
            response = await self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"}
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            print(f"Ошибка API: {e}")
            return None

    def create_pptx(self, content_json, design_name, filename="output.pptx"):
        """Создает .pptx файл с применением дизайна"""
        prs = Presentation()
        
        # Получаем цвета дизайна (по дефолту Basic Light)
        palette = self.design_palettes.get(design_name, self.design_palettes["Basic Light"])

        for slide_data in content_json['slides']:
            # Используем макет: Заголовок + Содержимое
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # --- ПРИМЕНЕНИЕ ДИЗАЙНА ---
            # Заливка фона
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*palette['bg'])

            # Настройка Заголовка
            title = slide.shapes.title
            title.text = slide_data['title']
            title_tf = title.text_frame
            title_p = title_tf.paragraphs[0]
            title_p
